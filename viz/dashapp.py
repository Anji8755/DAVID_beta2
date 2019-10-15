from dash.dependencies import Output, Input,State
import PyPDF2
import xlrd
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
import dash_core_components as dcc
import dash_html_components as html
from .server import server,app
from . import router
from django.shortcuts import redirect
from django.http import request
import base64
from urllib.parse import quote as urlquote
from tika import parser
import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy

import pandas as pd
import plotly.figure_factory as ff
import plotly.graph_objs as go
from django.db import connection
from django.contrib.auth.decorators import login_required
#specific to extracting information from word documents
import os
import subprocess
from docx import Document
import docx

import dash_table as dt

colors = {
    'background': '#C6B7B7',
    'text': '#7FDBFF'
}

app.css.append_css({
    "external_url": "https://codepen.io/chriddyp/pen/bWLwgP.css"
})

UPLOAD_DIRECTORY = os.getcwd()+"/dash-django-example/dash_test/viz/userUploads/"
SAVE_DIRECTORY = os.getcwd()+"/dash-django-example/dash_test/viz/CopiedFiles/"
TEXT_DIRECTORY= os.getcwd()+"/dash-django-example/dash_test/viz/TextFiles/"
BETA_TEXT_DIRECTORY= os.getcwd()+"/dash-django-example/dash_test/viz/textFilesBeta/"


app.layout = html.Div(style={'backgroundColor': "#fff"},children=[
    dcc.Location(id='url', refresh=False),
    #dcc.Link('Index', href=f'{app.url_base_pathname}'),
    ', ',
    html.Div(id='content'),
])


def fetch_data(q,returnNull):
    with connection.cursor() as cursor:
        cursor.execute(q)
        result={};
        if not returnNull:
            result = pd.read_sql(
            sql=q,
            con=connection
            )
        cursor = connection.cursor();

    return result


def fill_parseTab(FileID,orgFileName,UpldTime,UpldPath,UpldSize,GenFile,GenSize,GenPath,doctype,content,txt_file):
    content = content.replace("'", "''")
    query=(f''' Insert into parser_metadata (FileId,OriginalFile,UploadTime, UploadPath, UploadedSize,GeneratedFile,GeneratedSize,GeneratedPath,Doctype,File_Content,Txt_File) values('{FileID}','{orgFileName}','{UpldTime}','{UpldPath}','{UpldSize}','{GenFile}','{GenSize}','{GenPath}','{doctype}','{content}','{txt_file}') ''')
    print(query)
    pd.DataFrame(fetch_data(query,True))
    print("data inserted")
    return True


def save_file(name, content):
    """Decode and store a file uploaded with Plotly Dash."""
    #print(name)
    #print(content)
    data = content.encode("utf8").split(b";base64,")[1]
    print(os.path.join(UPLOAD_DIRECTORY, name))
    with open(os.path.join(UPLOAD_DIRECTORY, name), "wb") as fp:
        fp.write(base64.decodebytes(data))
    print("File saved on server ")


def uploaded_files():
    """List the files in the upload directory."""
    print("############inside uploaded files ##########################")
    files = []
    for filename in os.listdir(UPLOAD_DIRECTORY):
        path = os.path.join(UPLOAD_DIRECTORY, filename)
        if os.path.isfile(path):
            files.append(filename)
    #print(filename+str(datetime.datetime.now().timestamp()))
    return files

def parsed_files():
    """List the files in the upload directory."""
    files = []
    for filename in os.listdir(SAVE_DIRECTORY):
        path = os.path.join(SAVE_DIRECTORY, filename)
        if os.path.isfile(path):
            files.append(filename)
    return files

def extractTextfromCSV(filename):
    csv = pd.read_csv(filepath_or_buffer = filename)
    return True
def extractDataFromExcel(filename):
    print("inside excel module")
    os.chdir(UPLOAD_DIRECTORY)
    #loc = (filename)
    data=""
    # To open Workbook
    #wb = xlrd.open_workbook(filename)
    #sheet = wb.sheet_by_index(0)
    #for sheet in range(sheet.nrows):
    #    data+="\n----------------------------------NEW SHEET-------------------------------------\n"
    #    for row in range(sheet.nrows):
    #        data+="\n"
    #        print(sheet.row_values(row))
    #        rowvalue = sheet.row_values(row)
    #        for item in rowvalue:
    #            data+=item+"\t"
    #convert_to_txt(data,filename)
    #print("file generated for excel document")
    xl = pd.ExcelFile(filename)
    res = len(xl.sheet_names)
    print(xl)

    # For row 0 and column 0
    #sheet.cell_value(0, 0)


def file_download_link(filename,addPath):
    """Create a Plotly Dash 'A' element that downloads a file from the app."""
    print("/viz/{}/{}".format(addPath,urlquote(filename)))
    print("/viz/{}/{}".format(addPath,(filename)))
    location = "/viz/{}/{}".format(addPath,urlquote(filename))
    return html.A(filename, href=location)


#@app.callback(
#    Output("copied-file-list", "children"),
#    [Input("upload-data", "filename"), Input("upload-data", "contents")],
#)
#def renderParsedFiles(uploaded_filenames, uploaded_file_contents):
#    print("inside renderParsed files")
#    if uploaded_filenames is not None and uploaded_file_contents is not None:
#        for name, data in zip(uploaded_filenames, uploaded_file_contents):
#            save_file(name, data)
#            #parseData(name)

    #files = uploaded_files()
#    copiedFiles = parsed_files()
#    if len(copiedFiles) == 0:
#        print("inside renderParsed files1")
#        return [html.Li("No files yet!")]
#    else:
#        print("inside renderParsed files2")
#        return [html.Li(file_download_link(filename,"CopiedFiles")) for filename in copiedFiles]
def convert_to_txt(txt,filename):
    filename=filename+".txt"
    DIRECTORY = TEXT_DIRECTORY
    if "beta" in filename.lower() :
        DIRECTORY = BETA_TEXT_DIRECTORY
    with open(os.path.join(DIRECTORY, filename), "a+") as fp:
        print('inside with of convert_to_txt')
        #fp.write(base64.decodebytes(txt))
        fp.write(txt)
        print(fp)
    return

@app.callback(
    [Output("fileConnectorId", "value"),Output("copied-file-list", "children")],#,Output("formContainer", "style")
    [Input("upload-data", "filename"), Input("upload-data", "contents")],
)
def update_output(uploaded_filenames, uploaded_file_contents):
    print("inside function ##########vfgfgdfgdfg")
    """Save uploaded files and regenerate the file list."""
    #print(uploaded_filenames)
    fileList={}
            #return fileList[name]
    #print(uploaded_file_contents)
    if uploaded_filenames is not None and uploaded_file_contents is not None:
        for name, data in zip(uploaded_filenames, uploaded_file_contents):
            save_file(name, data)
            doctype=""
            content=""
            print(name)
            constname = str(datetime.datetime.now().timestamp())+name
            fileList[name]="Pending"
            isAllowed = True
            if 'xlsx' in name :
                    #doctype="xlsx"
                    #ret=extractDataFromExcel(name)
                    #setStat(ret)
                    isAllowed= False
                    fileList[name]="Success"
                    userAlert("Format not allowed " ,"Failure")
                    print("going for "+doctype)
            elif 'xls' in name :
                    #doctype="xls"
                    userAlert("Format not allowed " ,"Failure")
                    #ret=extractDataFromExcel(name)
                    isAllowed= False
                    fileList[name]="Failure"
                    #setStat(ret)
                    print("going for "+doctype)
            elif 'pptx' in name :
                    doctype="pptx"
                    ret=pptx(name,constname)
                    print("value from dictionary of pptx is: ",ret)
                    if ret['flag']==False:
                        #setStatus(ret)
                        fileList[name]="Success"
                        copy_slide_from_external_prs(name)
                    else :
                        fileList[name]="Failure"

            elif 'ppt' in name :
                    doctype="ppt"
                    ret=pptx(name,constname)
                    print("value from dictionary of pptx is: ",ret)
                    if ret['flag']==False:
                        #setStatus(ret)
                        fileList[name]="Success"
                        copy_slide_from_external_prs(name)
                    else:
                        fileList[name]="Failure"
            elif 'pdf' in name :
                    doctype="pdf"
                    ret=convert_pdf_to_txt(name,constname)
                    fileList[name]="Success"
                    #setStatus(ret)
            elif 'docx' in name :
                    doctype="docx"
                    fileList[name]="Success"
                    print('############about to enter parseData')
                    ret=parseData(name,constname)
                    if ret['flag']==True:
                        fileList[name]="Failure"
                    print("going for "+doctype)
            elif 'doc' in name :
                    doctype="doc"
                    print('############about to enter get_doc_text')
                    ret=parseData(name,constname)
                    fileList[name]="Success"
                    if ret['flag']==True:
                        fileList[name]="Failure"
                    #setStatus(ret)
                    print("going for "+doctype)
            elif 'csv' in name :
                    #doctype="csv"
                    userAlert("Format not allowed " ,"Failure")
                    isAllowed= False
                    #ret=extractTextfromCSV(name)
                    fileList[name]="Failure"
                    #setStatus(ret)
                    print("going for "+doctype)
            print("going to parse data")
            if isAllowed :
                Txt_File=constname
                fill_parseTab(constname,name,str(datetime.datetime.now().timestamp()),UPLOAD_DIRECTORY+name,"","OutputDoc"+name,"",SAVE_DIRECTORY+"OutputDoc"+name,doctype,content,Txt_File)
            #fileList[name]="success"
    files = uploaded_files()
    #print("inside renderParsed files100")
    if len(files) == 0:
        return "",""#,{"display":"none"}#[html.Li("No files yet!")]
    else:#[html.Li(file_download_link(filename,"userUploads")) for filename in files],

        return constname,[html.Li(file+" : "+stat) for file,stat in fileList.items()]#,{"display":"block"}


def pptx(filename,constname):
    print("hello i am u=inside pptx function")
    dict={}
    dict['flag']=False
    import pptx.exc
    os.chdir(UPLOAD_DIRECTORY)
    try:
        prs = Presentation(filename)
        print("----------------------")
        txt=""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
                    txt+=shape.text
        convert_to_txt(txt,constname)
        #return txt
        dict['txt']=txt
    except pptx.exc.PackageNotFoundError:#pptx.exc.PackageNotFoundError:
        print("this error pptx.exc.PackageNotFoundError")
        dict['flag']=True
        #dict['errMsg']='Erro occured due to incorrect file format. you may try *.pptx '
    print("dictionary from pptx is:",dict)
    return dict

def copy_slide_from_external_prs(filename):
    prs= Presentation()
    # copy from external presentation all objects into the existing presentationgetc
    os.chdir(UPLOAD_DIRECTORY)
    external_pres = Presentation(filename)
    print("copying ppt ")
    # specify the slide you want to copy the contents from
    #ext_slide = external_pres.slides[0]

    # Define the layout you want to use from your generated pptx
    SLD_LAYOUT = 5
    slide_layout = prs.slide_layouts[SLD_LAYOUT]

    # create now slide, to copy contents to
    curr_slide = prs.slides.add_slide(slide_layout)

    # now copy contents from external slide, but do not copy slide properties
    # e.g. slide layouts, etc., because these would produce errors, as diplicate
    # entries might be generated
    for ext_slide in external_pres.slides:
        curr_slide = prs.slides.add_slide(slide_layout)
        for shp in ext_slide.shapes:
            #if 'PICTURE' in shp.shape_type:
            #    pic = shp.image
            #    curr_slide.shapes._spTree.insert_element_before(pic, 'p:extLst')
            el = shp.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            print("shp is:",shp)
            print("shape type is: ",shp.shape_type)
            print("shape type is: ",type(shp.shape_type))
            #if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            #    print("shape type is: ",type(shp.shape_type))
            #    print("shape is:",shp)




    print("saving the ppt file now ")
    os.chdir(SAVE_DIRECTORY+"ppt/")
    print(os.getcwd())
    print("saving file"+filename)
    prs.save("outputPpt"+filename)
    return True

def pdf(filename):
        # creating a pdf file object
    os.chdir(UPLOAD_DIRECTORY)
    pdfFileObj = open(filename, 'rb')

    # creating a pdf reader object
    pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
    txt=""
    # printing number of pages in pdf file
    print(pdfReader.numPages)
    for page in range(pdfReader.numPages):
        pageObj = pdfReader.getPage(page)
        print(pageObj.extractText())
        text = pageObj.extractText()
                #print(text)
        txt+=text
        tokens = word_tokenize(text)
        #we'll create a new list which contains punctuation we wish to clean
        punctuations = ['(',')',';',':','[',']',',']
            #We initialize the stopwords variable which is a list of words like #"The", "I", "and", etc. that don't hold much value as keywords
        stop_words = stopwords.words('english')
        #We create a list comprehension which only returns a list of words #that are NOT IN stop_words and NOT IN punctuations.
        keywords = [word for word in tokens if not word in stop_words and not word in punctuations]
        print(keywords)
    #convert_to_txt(txt,constname)
    # creating a page object
    #pageObj = pdfReader.getPage(0)

    # extracting text from page


    # closing the pdf file object
    pdfFileObj.close()
    return True

def convert_pdf_to_txt(path_to_file,constname):
    print("inside convert_pdf_to_txt of pdfminer.six")
    os.chdir(UPLOAD_DIRECTORY)
    from io import StringIO
    from pdfminer.converter import TextConverter
    from pdfminer.pdfinterp import PDFPageInterpreter
    from pdfminer.pdfinterp import PDFResourceManager
    from pdfminer.pdfpage import PDFPage
    from pdfminer.layout import LAParams
    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    codec = 'utf-8'
    laparams = LAParams()
    device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)
    fp = open(path_to_file, 'rb')
    print(fp)
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    #password = ""
    maxpages = 0
    caching = True
   # pagenos=set()
    print("YOUr tesxtghdisdifdj;gofgjdlfg")
    #f = open("pdftxt.txt", "w")
    print(os.getcwd())
    txt=""
    for page in PDFPage.get_pages(fp):#, maxpages=maxpages, caching=caching, check_extractable=True):
        interpreter.process_page(page)

    text = retstr.getvalue()
        #print(text)
        #txt+=text
    convert_to_txt(text,constname)
    print('##################text inside pdf is:')
    print(text)
    #print('##################write content of admin log is')

    #f.close()

    fp.close()
    device.close()
    retstr.close()
    return text


def get_doc_text(filename,constname):
    #cmd = ['antiword', '-m', 'utf-8.txt', path]
    PIPE = subprocess.PIPE
    print("inside get_doc_text")
    path=UPLOAD_DIRECTORY+filename+"/"
    cmd = ['catdoc', '-d', 'utf-8', path]
    try:
        p = subprocess.Popen(cmd, stdout=PIPE)
        print("subprocess of p:",p)
        stdout, stderr = p.communicate()

        return stdout.decode('utf-8', 'ignore')
    except:
        return 'NOTHING HAPPENDED HERE'

def doc_to_text_catdoc(filename,constname):
    print("#####################inside doc_to_text_catdoc")
    os.chdir(UPLOAD_DIRECTORY)
    print('get file from:',UPLOAD_DIRECTORY+filename+"/")
    PIPE = subprocess.PIPE

    catdoc_cmd = ['catdoc', '-w' , name]
    catdoc_process = subprocess.Popen(catdoc_cmd, stdout=subprocess.PIPE,stderr=subprocess.PIPE)
    #for line in catdoc_process.stdout:
    #    if keyword in line:
    #        print line.strip()




    #p=subprocess.Popen('catdoc -w "%s"' % filename)
    p=subprocess.Popen(UPLOAD_DIRECTORY+filename)
    print('after subprocess p is:',p)
    #fi,fo,fe=p.stdin,p.stdout,p.stderr
    fi=p.stdin
    fo=p.stdout
    fe=p.stderr

    fi.close()
    txt = fo.read()
    print("txt is: ",txt)
    erroroutput = fe.read()
    fo.close()
    fe.close()
    if not erroroutput:
        return txt
    else:
        raise OSError("Executing the command caused an error: %s" % erroroutput)


def get_para_data(output_doc_name, paragraph,constname):
    """
    Write the run to the new file and then set its font, bold, alignment, color etc. data.
    """
    txt=""
    output_para = output_doc_name.add_paragraph()
    for run in paragraph.runs:
        output_run = output_para.add_run(run.text)
        print(output_run)
        txt+=str(output_run)
        print('documnet in txt format')
        print(txt)
        # Run's bold data
        output_run.bold = run.bold
        # Run's italic data
        output_run.italic = run.italic
        # Run's underline data
        output_run.underline = run.underline
        # Run's color data
        output_run.font.color.rgb = run.font.color.rgb
        # Run's font data
        output_run.style.name = run.style.name
    # Paragraph's alignment data
    convert_to_txt(txt,constname)
    output_para.paragraph_format.alignment = paragraph.paragraph_format.alignment
    return txt

def userAlert(msg,msgclass):
    print("Inside user alert function : " )
    val = js.call('userAlert', msg, msgClass)
    print (val)
    return True

def parseData(fileName,constname):
    #document = zipfile.ZipFile(UPLOAD_DIRECTORY+'fileName')
    #ZipFile.read(fileName, pwd=None)
        print("inside parse data")
        dict={}
        dict['flag']=False
        print(os.path.realpath('./'))
        txt=""
        os.chdir(UPLOAD_DIRECTORY)
        try:
            output_doc = Document()
            input_doc = Document(fileName)
            print(UPLOAD_DIRECTORY+fileName)
            for para in input_doc.paragraphs:
                txt+=get_para_data(output_doc, para,constname)
            os.chdir(SAVE_DIRECTORY+"doc/")
            print("current directory is ")
            dict['txt']=txt
        except KeyError :
            print("docx file format mismatch found ")
            dict['flag']=True
        print(os.getcwd())
        output_doc.save('OutputDoc'+fileName)
        return dict